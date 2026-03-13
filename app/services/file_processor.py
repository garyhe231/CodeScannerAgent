"""Extract text/content from uploaded files for use in chat."""
import base64
import io
from pathlib import Path
from typing import Dict

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
IMAGE_MEDIA_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}
AUDIO_VIDEO_EXTENSIONS = {".mp3", ".mp4", ".wav", ".m4a", ".mov", ".avi"}

MAX_FILE_BYTES = 20 * 1024 * 1024  # 20 MB


def process_file(filename: str, data: bytes) -> Dict:
    """
    Process uploaded file bytes and return a dict with:
      - filename
      - type: "text" | "image" | "unsupported"
      - content: extracted text (for text/doc types)
      - image_data: base64 string (for images)
      - media_type: MIME type (for images)
      - error: optional error message
    """
    suffix = Path(filename).suffix.lower()

    if suffix in IMAGE_EXTENSIONS:
        return _process_image(filename, data, suffix)
    elif suffix == ".pdf":
        return _process_pdf(filename, data)
    elif suffix in (".docx", ".doc"):
        return _process_docx(filename, data)
    elif suffix in (".xlsx", ".xls"):
        return _process_excel(filename, data)
    elif suffix in (".pptx", ".ppt"):
        return _process_pptx(filename, data)
    elif suffix in AUDIO_VIDEO_EXTENSIONS:
        return {
            "filename": filename,
            "type": "unsupported",
            "content": None,
            "image_data": None,
            "media_type": None,
            "error": f"Audio/video files ({suffix}) cannot be read — please provide a transcript instead.",
        }
    else:
        return _process_text(filename, data)


def _process_image(filename: str, data: bytes, suffix: str) -> Dict:
    media_type = IMAGE_MEDIA_TYPES.get(suffix, "image/jpeg")
    return {
        "filename": filename,
        "type": "image",
        "content": None,
        "image_data": base64.standard_b64encode(data).decode("utf-8"),
        "media_type": media_type,
        "error": None,
    }


def _process_text(filename: str, data: bytes) -> Dict:
    try:
        content = data.decode("utf-8", errors="replace")
        return {"filename": filename, "type": "text", "content": content, "image_data": None, "media_type": None, "error": None}
    except Exception as e:
        return {"filename": filename, "type": "text", "content": None, "image_data": None, "media_type": None, "error": str(e)}


def _process_pdf(filename: str, data: bytes) -> Dict:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[Page {i+1}]\n{text}")
            content = "\n\n".join(pages) if pages else "[No extractable text found in PDF]"
        return {"filename": filename, "type": "text", "content": content, "image_data": None, "media_type": None, "error": None}
    except Exception as e:
        return {"filename": filename, "type": "text", "content": None, "image_data": None, "media_type": None, "error": f"PDF read error: {e}"}


def _process_docx(filename: str, data: bytes) -> Dict:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return {"filename": filename, "type": "text", "content": content or "[Empty document]", "image_data": None, "media_type": None, "error": None}
    except Exception as e:
        return {"filename": filename, "type": "text", "content": None, "image_data": None, "media_type": None, "error": f"Word read error: {e}"}


def _process_excel(filename: str, data: bytes) -> Dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        content = "\n\n".join(sheets) if sheets else "[Empty spreadsheet]"
        return {"filename": filename, "type": "text", "content": content, "image_data": None, "media_type": None, "error": None}
    except Exception as e:
        return {"filename": filename, "type": "text", "content": None, "image_data": None, "media_type": None, "error": f"Excel read error: {e}"}


def _process_pptx(filename: str, data: bytes) -> Dict:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        content = "\n\n".join(slides) if slides else "[No text found in presentation]"
        return {"filename": filename, "type": "text", "content": content, "image_data": None, "media_type": None, "error": None}
    except Exception as e:
        return {"filename": filename, "type": "text", "content": None, "image_data": None, "media_type": None, "error": f"PowerPoint read error: {e}"}
